import os, io, datetime as dt
import streamlit as st
import requests, yfinance as yf, pandas as pd
from bs4 import BeautifulSoup
from pypdf import PdfReader
from pptx import Presentation
from openai import OpenAI
from psaw import PushshiftAPI
import feedparser

# ─── CONFIG & UI ───────────────────────────────────────────────────────────────
st.set_page_config(page_title="Smart Equity Dossier", layout="wide")
ACCENT = "#39ff14"
BG = "#191c1f"
CARD = "#2a2d31"
st.markdown(f"""
<style>
body {{background-color:{BG}; color:#e6eefa;}}
.stApp {{background:{BG};}}
.section-title {{color:{ACCENT}; font-size:1.25em; margin-top:1rem;}}
.card {{background:{CARD}; padding:1rem; border-radius:8px;}}
a {{color:{ACCENT};}}
</style>
""", unsafe_allow_html=True)

# ─── SECRETS & CLIENTS ─────────────────────────────────────────────────────────
PPLX_KEY = st.secrets.get("PPLX_API_KEY") or os.getenv("PPLX_API_KEY")
client = OpenAI(api_key=PPLX_KEY, base_url="https://api.perplexity.ai") if PPLX_KEY else None
UA = {"User-Agent":"Mozilla/5.0"}

# ─── HELPERS ────────────────────────────────────────────────────────────────────
def resolve_ticker(q):
    base = q.strip().upper()
    for suf in (".NS",".BO"):
        t = f"{base}{suf}"
        info = yf.Ticker(t).info
        if info.get("regularMarketPrice"):
            hist = yf.Ticker(t).history(period="5y")
            return t, info, hist
    return None, None, None

def fetch_rss(rss_url, n=5):
    resp = requests.get(rss_url, headers=UA, timeout=10)
    soup = BeautifulSoup(resp.content, "xml")
    items = soup.find_all("item")[:n]
    return [{"title":i.title.text, "link":i.link.text, "date":i.pubDate.text[:16]} for i in items]

def google_news(t, n=5):
    url = f"https://news.google.com/rss/search?q={t}+stock+india"
    return fetch_rss(url,n)

def et_news(rss, n=5):
    return fetch_rss(rss,n)

def mint_news(n=5):
    return fetch_rss("https://www.livemint.com/rss/news", n)

def scrape_screener(code, idx):
    url = f"https://www.screener.in/company/{code}/consolidated/"
    r = requests.get(url, headers=UA, timeout=10)
    soup = BeautifulSoup(r.text,"html.parser")
    tbls = soup.find_all("table",class_="data-table")
    if len(tbls)>idx:
        return pd.read_html(str(tbls[idx]))[0]
    return None

def fetch_vp(n=15):
    feed = feedparser.parse("https://valuepickr4.rssing.com/chan-72344682/latest.php")
    cutoff = dt.datetime.utcnow() - dt.timedelta(days=90)
    threads=[]
    for e in feed.entries:
        pub=dt.datetime(*e.published_parsed[:6])
        if pub<cutoff: continue
        import re
        m=re.search(r"\[Replies:\s*(\d+)\]",e.title)
        cnt=int(m.group(1)) if m else 0
        threads.append({"date":pub,"replies":cnt,"title":e.title.split(" [")[0],"link":e.link})
    df=pd.DataFrame(threads)
    return df.sort_values(["replies","date"],ascending=False).head(n)

def fetch_reddit(n=15):
    api=PushshiftAPI()
    now=int(dt.datetime.utcnow().timestamp()); since=now-90*24*3600
    rec=[]
    for sub in ["IndianStockMarket","Stocks","investingindia"]:
        for p in api.search_submissions(subreddit=sub, after=since, before=now,
                                        filter=["title","num_comments","score","created_utc","full_link"],
                                        sort="desc",sort_type="num_comments", limit=n):
            rec.append({"date":dt.datetime.utcfromtimestamp(p.created_utc),
                        "sub":sub,"comments":p.num_comments,"score":p.score,
                        "title":p.title,"link":p.full_link})
    return pd.DataFrame(rec).sort_values(["comments","score"],ascending=False).head(n)

def extract_text(b, fn):
    if fn.lower().endswith(".pdf"):
        return "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(b)).pages)
    prs=Presentation(io.BytesIO(b))
    txt=[shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape,"text")]
    return "\n".join(txt)

def analyze_doc(text,prompt):
    if not client: return "🔑 Missing Perplexity key."
    msgs=[{"role":"system","content":prompt},{"role":"user","content":text[:12000]}]
    res=client.chat.completions.create(model="sonar-pro",messages=msgs)
    return res.choices[0].message.content

# ─── SIDEBAR ───────────────────────────────────────────────────────────────────
st.sidebar.title("🔍 Search & Docs")
query=st.sidebar.text_input("Ticker/Company","")
uploads=st.sidebar.file_uploader("📂 Upload PDF/PPTX",accept_multiple_files=True,type=["pdf","pptx"])

# ─── MAIN ──────────────────────────────────────────────────────────────────────
if not query:
    st.info("Enter a ticker in the sidebar to begin.")
    st.stop()

ticker, info, hist = resolve_ticker(query)
if not ticker:
    st.error("Ticker not found on Yahoo Finance; try again.")
    st.stop()

code=ticker.replace(".NS","").replace(".BO","")
st.markdown(f"<div class='section-title'>💹 {ticker}</div>",unsafe_allow_html=True)

# Price & Chart
c1,c2,c3=st.columns(3)
c1.metric("Price",f"₹{info['regularMarketPrice']}")
c2.metric("52W High",f"₹{info['fiftyTwoWeekHigh']}")
c3.metric("52W Low",f"₹{info['fiftyTwoWeekLow']}")
st.line_chart(hist["Close"])

# 1 Sector Trends & Triggers
st.markdown("<div class='section-title'>🌐 Sectoral Trends & Triggers</div>",unsafe_allow_html=True)
for n in google_news(info.get("sector",""),5):
    st.write(f"- [{n['title']}]({n['link']}) — {n['date']}")

# 2 News & Competition
st.markdown("<div class='section-title'>📰 News & Competition</div>",unsafe_allow_html=True)
st.markdown("**Google News**")
for n in google_news(code,5): st.write(f"- [{n['title']}]({n['link']}) — {n['date']}")
st.markdown("**ET Top Stories**")
for n in et_news("https://economictimes.indiatimes.com/rssfeedstopstories.cms",5):
    st.write(f"- [{n['title']}]({n['link']}) — {n['date']}")
st.markdown("**ET Industry News**")
for n in et_news("https://economictimes.indiatimes.com/rss/etindustryrss.cms",5):
    st.write(f"- [{n['title']}]({n['link']}) — {n['date']}")
st.markdown("**Mint News**")
for n in mint_news(5): st.write(f"- [{n['title']}]({n['link']}) — {n['date']}")

# 6 Forum Threads
st.markdown("<div class='section-title'>💬 Community Threads (90d)</div>",unsafe_allow_html=True)
st.markdown("**ValuePickr**")
vp=fetch_vp(15)
st.dataframe(vp.rename(columns={"date":"Date","replies":"Replies","title":"Topic","link":"URL"}))
st.markdown("**Reddit**")
rd=fetch_reddit(15)
st.dataframe(rd.rename(columns={"date":"Date","sub":"Subreddit","comments":"Comments","score":"Score","title":"Title","link":"URL"}))

# Perplexity AI Sections
if uploads:
    st.markdown("<div class='section-title'>🤖 AI-Powered Management Analysis</div>",unsafe_allow_html=True)
    for f in uploads:
        b, name = f.read(), f.name
        txt=extract_text(b,name)
        with st.spinner(f"Analyzing {name}…"):
            summ=analyze_doc(txt,"Extract management commentary: strategy, guidance, risks.")
            integ=analyze_doc(txt,"Score integrity on Guidance, Delivery, Transparency, Governance.")
        st.subheader(name)
        st.write(summ)
        st.write("**Integrity Matrix:**", integ)
